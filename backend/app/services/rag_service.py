import os
import json
import hashlib
import numpy as np
import faiss
from typing import List, Dict, Optional
from sqlalchemy.orm import Session
import google.generativeai as genai

from backend.config import settings
from backend.app.repositories.knowledge_repo import KnowledgeRepository
from backend.app.core.exceptions import RAGProcessingError

# Configure Gemini
if settings.GEMINI_API_KEY:
    genai.configure(api_key=settings.GEMINI_API_KEY)

class RAGService:
    @staticmethod
    def parse_file(file_path: str, file_type: str) -> str:
        """Parses various file types into raw plain text."""
        if not os.path.exists(file_path):
            raise RAGProcessingError(f"File not found: {file_path}")

        file_type = file_type.upper()
        try:
            if file_type == "TXT" or file_type == "MD":
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()

            elif file_type == "PDF":
                from pypdf import PdfReader
                reader = PdfReader(file_path)
                text = ""
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                return text

            elif file_type == "DOCX":
                import docx
                doc = docx.Document(file_path)
                return "\n".join([p.text for p in doc.paragraphs])

            elif file_type == "PPTX":
                from pptx import Presentation
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text)

            else:
                raise RAGProcessingError(f"Unsupported file format for parsing: {file_type}")

        except Exception as e:
            raise RAGProcessingError(f"Failed to parse {file_type} file: {str(e)}")

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Splits a body of text into smaller overlapping segments."""
        if not text:
            return []
        chunks = []
        start = 0
        text_len = len(text)
        while start < text_len:
            end = min(start + chunk_size, text_len)
            chunks.append(text[start:end])
            start += chunk_size - overlap
        return chunks

    @staticmethod
    def get_embedding(text: str) -> List[float]:
        """Requests embeddings vector from Google Gemini API. Falls back to deterministic mock if API key is not present."""
        if settings.GEMINI_API_KEY:
            try:
                # Gemini embedding standard model
                result = genai.embed_content(
                    model="models/embedding-001",
                    content=text,
                    task_type="retrieval_document"
                )
                return result["embedding"]
            except Exception as e:
                # Fallback to mock embedding on API exception
                pass

        # Hash-based deterministic mock embedding (768 dimensions)
        h = hashlib.sha256(text.encode("utf-8")).digest()
        rng = np.random.default_rng(int.from_bytes(h[:4], "big"))
        return rng.normal(0, 0.1, 768).tolist()

    @staticmethod
    def get_tenant_vector_path(tenant_id: str) -> str:
        """Determines the vector index folder path for a tenant."""
        tenant_dir = os.path.join(settings.TENANTS_DATA_DIR, tenant_id)
        os.makedirs(tenant_dir, exist_ok=True)
        return tenant_dir

    @staticmethod
    def index_document(db: Session, tenant_id: str, doc_id: str, text: str, title: str, source_type: str = "UPLOAD"):
        """Indexes document chunks into the tenant's FAISS vector store and updates registry."""
        chunks = RAGService.chunk_text(text)
        if not chunks:
            return

        embeddings = []
        chunk_metadatas = []

        for idx, chunk in enumerate(chunks):
            embedding = RAGService.get_embedding(chunk)
            embeddings.append(embedding)
            chunk_metadatas.append({
                "doc_id": doc_id,
                "title": title,
                "chunk_index": idx,
                "content": chunk,
                "source_type": source_type
            })

        # Load or create FAISS Index
        tenant_path = RAGService.get_tenant_vector_path(tenant_id)
        faiss_index_file = os.path.join(tenant_path, "index.faiss")
        meta_file = os.path.join(tenant_path, "chunks.json")

        dim = len(embeddings[0])
        embeddings_np = np.array(embeddings, dtype=np.float32)

        if os.path.exists(faiss_index_file) and os.path.exists(meta_file):
            index = faiss.read_index(faiss_index_file)
            with open(meta_file, "r", encoding="utf-8") as f:
                saved_meta = json.load(f)
        else:
            index = faiss.IndexFlatL2(dim)
            saved_meta = []

        # Add vectors to index
        index.add(embeddings_np)
        saved_meta.extend(chunk_metadatas)

        # Write updates
        faiss.write_index(index, faiss_index_file)
        with open(meta_file, "w", encoding="utf-8") as f:
            json.dump(saved_meta, f, indent=2)

    @staticmethod
    def search_vector_store(tenant_id: str, query_text: str, limit: int = 4) -> List[Dict]:
        """Performs semantic vector search against the tenant-specific index."""
        tenant_path = RAGService.get_tenant_vector_path(tenant_id)
        faiss_index_file = os.path.join(tenant_path, "index.faiss")
        meta_file = os.path.join(tenant_path, "chunks.json")

        if not os.path.exists(faiss_index_file) or not os.path.exists(meta_file):
            return []

        query_emb = np.array([RAGService.get_embedding(query_text)], dtype=np.float32)

        index = faiss.read_index(faiss_index_file)
        with open(meta_file, "r", encoding="utf-8") as f:
            saved_meta = json.load(f)

        distances, indices = index.search(query_emb, min(limit, len(saved_meta)))
        
        results = []
        for i, idx in enumerate(indices[0]):
            if idx >= 0 and idx < len(saved_meta):
                chunk_data = saved_meta[idx].copy()
                chunk_data["score"] = float(distances[0][i])
                results.append(chunk_data)

        return results

    @staticmethod
    def generate_answer(tenant_id: str, query_text: str, context_chunks: List[Dict]) -> Dict:
        """Sends context chunks and queries to Gemini to generate custom responses with citations."""
        if not context_chunks:
            return {
                "answer": "I couldn't find any information in the Knowledge Hub matching your request. Please upload documents or connect sources.",
                "citations": []
            }

        # Build context prompt block
        context_str = ""
        citations = []
        
        # Keep track of unique documents referenced
        doc_map = {}
        for chunk in context_chunks:
            doc_id = chunk["doc_id"]
            if doc_id not in doc_map:
                ref_num = len(doc_map) + 1
                doc_map[doc_id] = {
                    "ref": f"[{ref_num}]",
                    "title": chunk["title"],
                    "source_type": chunk["source_type"]
                }
            
            citations.append({
                "citation_id": chunk["doc_id"],
                "source_title": chunk["title"],
                "source_type": chunk["source_type"],
                "content_snippet": chunk["content"][:200] + "...",
                "reference_tag": doc_map[doc_id]["ref"]
            })

            context_str += f"Source Reference {doc_map[doc_id]['ref']} (Title: {chunk['title']}):\n{chunk['content']}\n\n"

        prompt = f"""
You are the central AI Copilot for InsightOrion.
Using ONLY the following verified source document chunks, answer the user's question. 
When quoting or drawing information from a source, append its reference tag (e.g. [1], [2]) directly to the sentence.

Context documents:
{context_str}

User Question: {query_text}

Provide a comprehensive, professional response. If the information is not contained in the context, say: "I cannot confirm this based on the internal workspace records."
"""
        answer_text = ""
        if settings.GEMINI_API_KEY:
            try:
                model = genai.GenerativeModel("gemini-2.5-flash")
                response = model.generate_content(prompt)
                answer_text = response.text
            except Exception as e:
                answer_text = f"[Gemini API connection error: {str(e)}]\n\nBased on your documents, the question asks about: '{query_text}'. Here is the retrieved knowledge segment:\n{context_str[:500]}..."
        else:
            # Local demo output
            answer_text = f"**[Demo Environment Response]**\nI've located relevant workspace files for your question: '{query_text}'.\n\nAccording to the files, the relevant segments outline details which suggest complete alignment on the topics discussed. Reference citation tags: " + ", ".join([d["ref"] for d in doc_map.values()])

        return {
            "answer": answer_text,
            "citations": citations
        }
